from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
FIGURES_DIR = BASE_DIR / "figures"
RESULTS_CSV = BASE_DIR / "results.csv"
OUTPUT_PPTX = BASE_DIR / "ML-KEM_Performance_Analysis_Presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(27, 45, 71)
BLUE = RGBColor(64, 110, 155)
TEAL = RGBColor(86, 139, 133)
AMBER = RGBColor(188, 141, 88)
ROSE = RGBColor(170, 109, 101)
INK = RGBColor(40, 52, 66)
MUTED = RGBColor(98, 109, 122)
BG = RGBColor(249, 250, 252)
GRID = RGBColor(223, 229, 238)
WHITE = RGBColor(255, 255, 255)
PALE_BLUE = RGBColor(236, 242, 250)
PALE_GREEN = RGBColor(233, 245, 241)
PALE_GOLD = RGBColor(249, 239, 224)
PALE_ROSE = RGBColor(246, 234, 232)

FONT_CN = "PingFang SC"
FONT_EN = "Aptos"


def load_data() -> pd.DataFrame:
    df = pd.read_csv(RESULTS_CSV)
    df["online_ms"] = df["encaps_mean_ms"] + df["decaps_mean_ms"]
    df["full_ms"] = df["keygen_mean_ms"] + df["online_ms"]
    df["full_ops"] = 1000 / df["full_ms"]
    df["online_ops"] = 1000 / df["online_ms"]
    baseline = df.iloc[0]
    df["keygen_norm"] = df["keygen_mean_ms"] / baseline["keygen_mean_ms"]
    df["encaps_norm"] = df["encaps_mean_ms"] / baseline["encaps_mean_ms"]
    df["decaps_norm"] = df["decaps_mean_ms"] / baseline["decaps_mean_ms"]
    df["online_norm"] = df["online_ms"] / baseline["online_ms"]
    df["full_norm"] = df["full_ms"] / baseline["full_ms"]
    df["pk_ct_bytes"] = df["public_key_bytes"] + df["ciphertext_bytes"]
    return df


def set_font(run, size: int, bold: bool = False, color: RGBColor = INK, font_name: str = FONT_CN) -> None:
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.35), Inches(12.0), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    set_font(run, 24, bold=True, color=NAVY)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.68), Inches(0.92), Inches(11.8), Inches(0.4))
        p2 = sub_box.text_frame.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        set_font(run2, 10, color=MUTED)


def add_footer(slide, page_num: int) -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(7.08), Inches(12.0), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = GRID
    line.line.fill.background()
    foot = slide.shapes.add_textbox(Inches(11.95), Inches(7.11), Inches(0.4), Inches(0.2))
    p = foot.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(page_num)
    set_font(run, 9, color=MUTED, font_name=FONT_EN)


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, size: int = 16) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.bullet = True
        p.level = 0
        p.space_after = Pt(8)
        set_font(p.runs[0], size, color=INK)


def add_note_box(slide, title: str, body: str, left: float, top: float, width: float, height: float, fill_color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    set_font(r1, 13, bold=True, color=NAVY)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = body
    set_font(r2, 11, color=INK)


def add_image(slide, image_name: str, left: float, top: float, width: float | None = None, height: float | None = None) -> None:
    image_path = FIGURES_DIR / image_name
    kwargs = {}
    if width is not None:
        kwargs["width"] = Inches(width)
    if height is not None:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), **kwargs)


def add_table(slide, data: list[list[str]], columns: list[str], left: float, top: float, width: float, height: float) -> None:
    rows = len(data) + 1
    cols = len(columns)
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    col_width = Inches(width / cols)
    for c in range(cols):
        table.columns[c].width = col_width
        cell = table.cell(0, c)
        cell.text = columns[c]
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                set_font(run, 11, bold=True, color=WHITE)
    for r, row_data in enumerate(data, start=1):
        for c, value in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else PALE_BLUE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    set_font(run, 11, bold=(c == 0), color=INK)


def add_metric_card(slide, title: str, value: str, caption: str, left: float, top: float, width: float, accent: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = GRID
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(0.08), Inches(1.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Pt(14)
    tf.margin_top = Pt(10)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    set_font(r1, 11, bold=True, color=MUTED)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    set_font(r2, 20, bold=True, color=NAVY, font_name=FONT_EN)
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = caption
    set_font(r3, 10, color=INK)


def build_presentation() -> Path:
    df = load_data()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.35))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    hero = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.5), Inches(2.0))
    p = hero.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "ML-KEM-512/768/1024\n性能比较实验与安全代价分析"
    set_font(run, 26, bold=True, color=NAVY)
    sub = hero.text_frame.add_paragraph()
    sub_run = sub.add_run()
    sub_run.text = "基于 liboqs-python 的实验评测与部署取舍讨论"
    set_font(sub_run, 15, color=MUTED)
    add_note_box(
        slide,
        "核心问题",
        "在相同实验环境下，安全等级上升到底会带来多少时间、空间与系统吞吐代价？",
        0.82,
        4.2,
        5.2,
        1.2,
        PALE_BLUE,
    )
    add_image(slide, "visual_pareto_bubble.png", 6.8, 1.55, width=5.55)
    footer = slide.shapes.add_textbox(Inches(0.85), Inches(6.7), Inches(5.0), Inches(0.3))
    p = footer.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "项目输出：基准测试、归一化代价、整体握手成本、吞吐能力、三维权衡图"
    set_font(r, 10, color=MUTED)

    # 2. roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "汇报路线", "从“做了什么”到“该怎么选型”")
    steps = [
        "研究背景：为什么比较 ML-KEM-512 / 768 / 1024",
        "实验设计：测试环境、测量方法、指标定义",
        "结果一：单步运行时间与原始数据快照",
        "结果二：归一化代价，安全升级到底贵了多少",
        "结果三：完整密钥建立成本与吞吐能力",
        "结果四：速度-尺寸-安全性的综合权衡",
        "结论：不同部署目标下的参数集建议",
    ]
    add_bullets(slide, steps, 0.9, 1.5, 7.2, 4.8, size=18)
    add_note_box(
        slide,
        "一句话判断",
        "ML-KEM-512 最快，ML-KEM-1024 最强但代价最高，ML-KEM-768 是最均衡的折中点。",
        8.5,
        2.1,
        3.8,
        1.5,
        PALE_GREEN,
    )
    add_metric_card(slide, "测量轮数", "500", "每个参数集均重复测量", 8.55, 4.2, 1.7, TEAL)
    add_metric_card(slide, "在线成本", "Encaps + Decaps", "模拟实际协商阶段", 10.35, 4.2, 1.8, BLUE)
    add_metric_card(slide, "完整成本", "KeyGen + Online", "纳入密钥生成开销", 8.55, 5.8, 3.6, AMBER)
    add_footer(slide, 2)

    # 3. background
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "研究对象与评价维度", "同一标准下，不同参数集体现出不同的安全-性能取舍")
    add_note_box(slide, "参数集", "ML-KEM-512、ML-KEM-768、ML-KEM-1024", 0.8, 1.4, 3.45, 1.1, PALE_BLUE)
    add_note_box(slide, "安全等级", "分别对应 Level 1 / 3 / 5 的递增安全强度", 4.5, 1.4, 3.55, 1.1, PALE_GREEN)
    add_note_box(slide, "比较目标", "不仅看谁更快，还要看安全升级带来的额外代价", 8.2, 1.4, 4.0, 1.1, PALE_GOLD)
    items = [
        "时间代价：KeyGen、Encaps、Decaps 的平均时延",
        "空间代价：公钥、私钥、密文长度",
        "流程代价：一次在线协商与完整密钥建立的总成本",
        "系统能力：单位时间可处理多少次密钥建立请求",
    ]
    add_bullets(slide, items, 0.95, 3.0, 6.5, 3.0, size=17)
    add_table(
        slide,
        [
            ["ML-KEM-512", "L1", "800", "1632", "768"],
            ["ML-KEM-768", "L3", "1184", "2400", "1088"],
            ["ML-KEM-1024", "L5", "1568", "3168", "1568"],
        ],
        ["参数集", "安全级别", "公钥(B)", "私钥(B)", "密文(B)"],
        7.4,
        3.0,
        5.2,
        2.5,
    )
    add_footer(slide, 3)

    # 4. setup
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "实验环境与方法", "保证三个参数集在同一软件栈和同一测试流程下比较")
    left_items = [
        "实现库：liboqs-python",
        "算法对象：ML-KEM-512 / 768 / 1024",
        "预热轮数：20",
        "正式测量轮数：500",
        "正确性校验：验证封装/解封装后的共享密钥一致",
    ]
    right_items = [
        "测量操作：KeyGen、Encaps、Decaps",
        "统计输出：平均值与标准差",
        "衍生指标：online total、full total、throughput",
        "归一化基准：ML-KEM-512 = 1.00",
    ]
    add_bullets(slide, left_items, 0.9, 1.5, 5.6, 3.8, size=16)
    add_bullets(slide, right_items, 6.7, 1.5, 5.6, 3.4, size=16)
    add_note_box(
        slide,
        "方法优势",
        "把单步性能、整体握手成本与系统容量三种视角合在一起，结论更接近真实部署场景。",
        0.9,
        5.55,
        11.4,
        1.0,
        PALE_ROSE,
    )
    add_footer(slide, 4)

    # 5. raw results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "原始结果快照", "先看绝对时间：参数集越大，三个核心操作都变慢")
    raw_rows = []
    for _, row in df.iterrows():
        raw_rows.append(
            [
                row["algorithm"],
                str(int(row["security_level"])),
                f'{row["keygen_mean_ms"]:.4f}',
                f'{row["encaps_mean_ms"]:.4f}',
                f'{row["decaps_mean_ms"]:.4f}',
                f'{row["public_key_bytes"]:.0f}',
                f'{row["ciphertext_bytes"]:.0f}',
            ]
        )
    add_table(
        slide,
        raw_rows,
        ["参数集", "级别", "KeyGen (ms)", "Encaps (ms)", "Decaps (ms)", "公钥(B)", "密文(B)"],
        0.8,
        1.55,
        7.4,
        2.3,
    )
    add_image(slide, "keygen_time.png", 8.55, 1.45, width=3.65)
    add_image(slide, "encaps_time.png", 8.45, 3.45, width=1.82)
    add_image(slide, "decaps_time.png", 10.35, 3.45, width=1.82)
    add_note_box(
        slide,
        "原始观察",
        "从 512 升级到 1024 后，三类操作的平均时延都显著上升；这意味着更高安全等级并不是“免费升级”。",
        0.86,
        4.45,
        7.2,
        1.55,
        PALE_GOLD,
    )
    add_footer(slide, 5)

    # 6. normalized
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "归一化代价分析", "以 ML-KEM-512 为基准，看安全升级到底贵了多少")
    add_image(slide, "visual_heatmap_tradeoffs.png", 0.75, 1.4, width=8.0)
    add_note_box(slide, "关键信息 1", "ML-KEM-768 在时间和尺寸上大致是 1.4x 到 1.6x。", 8.95, 1.7, 3.0, 1.0, PALE_GREEN)
    add_note_box(slide, "关键信息 2", "ML-KEM-1024 在时间上达到约 2.1x 到 2.5x，密钥与密文尺寸接近 2x。", 8.95, 2.9, 3.0, 1.25, PALE_GOLD)
    add_note_box(slide, "结论", "从“增量代价”看，768 是合理折中；1024 的高安全带来非常明确的额外成本。", 8.95, 4.45, 3.0, 1.15, PALE_ROSE)
    add_metric_card(slide, "KeyGen 放大", f'{df.loc[2, "keygen_norm"]:.2f}x', "1024 相对 512", 8.95, 5.9, 1.35, BLUE)
    add_metric_card(slide, "Online 放大", f'{df.loc[2, "online_norm"]:.2f}x', "1024 相对 512", 10.45, 5.9, 1.35, AMBER)
    add_footer(slide, 6)

    # 7. workflow latency
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "完整密钥建立成本", "相比单步时间，更需要关注一次完整协商到底要付出多少")
    add_image(slide, "visual_latency_dumbbell.png", 0.72, 1.35, width=8.25)
    add_metric_card(slide, "Online Total", f'{df.loc[0, "online_ms"]:.4f} ms', "ML-KEM-512", 9.15, 1.75, 1.45, TEAL)
    add_metric_card(slide, "Online Total", f'{df.loc[1, "online_ms"]:.4f} ms', "ML-KEM-768", 10.75, 1.75, 1.45, AMBER)
    add_metric_card(slide, "Online Total", f'{df.loc[2, "online_ms"]:.4f} ms', "ML-KEM-1024", 9.15, 3.45, 1.45, ROSE)
    add_metric_card(slide, "Full Total", f'{df.loc[2, "full_ms"]:.4f} ms', "1024 最完整成本", 10.75, 3.45, 1.45, BLUE)
    add_note_box(
        slide,
        "部署理解",
        "如果系统需要频繁建立会话密钥，那么 full total 比单步均值更有解释力；768 的整体成本约为 512 的 1.55x，1024 约为 2.26x。",
        9.15,
        5.25,
        3.1,
        1.15,
        PALE_BLUE,
    )
    add_footer(slide, 7)

    # 8. throughput
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "吞吐能力分析", "更高安全等级不仅让单次操作更慢，也会降低系统每秒可处理请求数")
    add_image(slide, "paper_table_throughput.png", 0.75, 1.4, width=8.0)
    add_note_box(slide, "系统容量", "完整密钥建立吞吐从 36,976 ops/s 下降到 16,380 ops/s，1024 仅约为 512 的 44%。", 9.1, 1.85, 3.1, 1.2, PALE_GOLD)
    add_note_box(slide, "实际意义", "在高并发服务器场景下，这种吞吐下降会直接反映为更高 CPU 压力与更低请求处理能力。", 9.1, 3.35, 3.1, 1.15, PALE_GREEN)
    add_metric_card(slide, "Full Throughput", f'{df.loc[0, "full_ops"]:.0f} ops/s', "ML-KEM-512", 9.1, 5.0, 1.45, TEAL)
    add_metric_card(slide, "Full Throughput", f'{df.loc[1, "full_ops"]:.0f} ops/s', "ML-KEM-768", 10.7, 5.0, 1.45, AMBER)
    add_metric_card(slide, "Full Throughput", f'{df.loc[2, "full_ops"]:.0f} ops/s', "ML-KEM-1024", 9.9, 6.15, 1.45, ROSE)
    add_footer(slide, 8)

    # 9. overall tradeoff
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "三维权衡：速度、尺寸与安全", "把结论放回真实决策：不是单看快慢，而是看“值不值得”")
    add_image(slide, "visual_pareto_bubble.png", 0.7, 1.2, width=8.4)
    add_note_box(slide, "ML-KEM-512", "速度最快、尺寸最小，适合性能最敏感的环境。", 9.2, 1.55, 3.0, 1.0, PALE_GREEN)
    add_note_box(slide, "ML-KEM-768", "整体最均衡，是安全与代价之间最像“默认选择”的参数集。", 9.2, 2.8, 3.0, 1.0, PALE_GOLD)
    add_note_box(slide, "ML-KEM-1024", "安全最强，但时延、尺寸和吞吐代价都最重，适合高安全优先级场景。", 9.2, 4.05, 3.0, 1.1, PALE_ROSE)
    add_footer(slide, 9)

    # 10. conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "结论与选型建议", "把实验结果转成可以直接表达的汇报结论")
    add_note_box(slide, "结论 1", "ML-KEM 参数集越高，计算时间、数据规模和系统吞吐损失越明显。", 0.85, 1.45, 11.5, 0.9, PALE_BLUE)
    add_note_box(slide, "结论 2", "ML-KEM-768 在安全增强与性能成本之间表现出最清晰的折中属性。", 0.85, 2.55, 11.5, 0.9, PALE_GREEN)
    add_note_box(slide, "结论 3", "ML-KEM-1024 更适合安全优先的关键场景，而不是默认全局替换。", 0.85, 3.65, 11.5, 0.9, PALE_ROSE)
    add_metric_card(slide, "性能优先", "ML-KEM-512", "低时延 / 高吞吐", 1.0, 5.05, 2.6, TEAL)
    add_metric_card(slide, "均衡默认", "ML-KEM-768", "综合表现最佳", 5.0, 5.05, 2.6, AMBER)
    add_metric_card(slide, "安全优先", "ML-KEM-1024", "更高保护等级", 9.0, 5.05, 2.6, ROSE)
    add_footer(slide, 10)

    prs.save(OUTPUT_PPTX)
    return OUTPUT_PPTX


if __name__ == "__main__":
    output = build_presentation()
    print(output)
